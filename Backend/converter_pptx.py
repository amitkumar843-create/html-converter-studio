import uuid
import tempfile
import math
import re
from pathlib import Path

from bs4 import BeautifulSoup
from playwright.async_api import async_playwright

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE_DIR = Path(__file__).parent
OUTPUT_DIR = BASE_DIR / "output"
OUTPUT_DIR.mkdir(exist_ok=True)



def safe_pptx_filename(source_filename: str | None = None) -> str:
    """Create PPTX output name from original uploaded HTML file name."""
    if not source_filename:
        return f"{uuid.uuid4().hex}.pptx"

    stem = Path(str(source_filename)).stem.strip()
    if not stem:
        return f"{uuid.uuid4().hex}.pptx"

    # Windows-safe filename cleanup.
    stem = re.sub(r'[<>:"/\\|?*\x00-\x1F]', '_', stem)
    stem = re.sub(r'\s+', ' ', stem).strip().rstrip('.')
    if not stem:
        return f"{uuid.uuid4().hex}.pptx"

    return f"{stem}.pptx"


def unique_output_path(file_name: str) -> Path:
    """Avoid overwriting existing PPTX by appending _1, _2, etc."""
    output_path = OUTPUT_DIR / file_name
    if not output_path.exists():
        return output_path

    stem = Path(file_name).stem
    suffix = Path(file_name).suffix or ".pptx"
    counter = 1
    while True:
        candidate = OUTPUT_DIR / f"{stem}_{counter}{suffix}"
        if not candidate.exists():
            return candidate
        counter += 1

def is_slide_deck(html_content: str) -> bool:
    soup = BeautifulSoup(html_content, "html.parser")
    return len(soup.select(".slide")) > 1

def count_slides(html_content: str) -> int:
    soup = BeautifulSoup(html_content, "html.parser")
    # Support normal decks (.slide) and wrapper/iframe decks such as GNIDA_v5 (.slide-shell)
    count = len(soup.select(".slide"))
    if count == 0:
        count = len(soup.select(".slide-shell"))
    # FIX: If no slide container exists, treat the entire body as 1 slide
    return count if count > 0 else 1

def build_layered_pptx(slide_data_list, output_file):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    
    for idx, slide_data in enumerate(slide_data_list):
        if idx == 0:
            prs.slide_width = Inches(slide_data["width"] / 96.0)
            prs.slide_height = Inches(slide_data["height"] / 96.0)
        
        slide = prs.slides.add_slide(blank_slide_layout)
        clip_x = slide_data['clip_x']
        clip_y = slide_data['clip_y']

        slide.shapes.add_picture(
            str(slide_data['bg_img_path']), 0, 0,
            width=Inches(slide_data["width"] / 96.0), 
            height=Inches(slide_data["height"] / 96.0)
        )
        
        for comp in slide_data['components']:
            cx = Inches((comp['x'] - clip_x) / 96.0)
            cy = Inches((comp['y'] - clip_y) / 96.0)
            cw = Inches(comp['w'] / 96.0)
            ch = Inches(comp['h'] / 96.0)
            try:
                slide.shapes.add_picture(str(comp['img_path']), cx, cy, width=cw, height=ch)
            except Exception:
                pass

        for el in slide_data['elements']:
            x_px = el['x'] - clip_x
            y_px = el['y'] - clip_y
            
            if x_px < -10 or y_px < -10: continue
                
            x = Inches(max(0, x_px) / 96.0)
            y = Inches(max(0, y_px) / 96.0)
            w = Inches(max(1, el['w']) / 96.0)
            h = Inches(max(1, el['h']) / 96.0)
            
            try:
                txBox = slide.shapes.add_textbox(x, y, w, h)
                tf = txBox.text_frame
                tf.clear()
                tf.margin_left = 0
                tf.margin_top = 0
                tf.margin_right = 0
                tf.margin_bottom = 0
                tf.word_wrap = True
                
                if el['vAlign'] == 'middle':
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                p = tf.paragraphs[0]
                p.text = el['text']
                p.font.size = Pt(max(1, el['fontSize'] * 0.75))
                
                if el['color']:
                    p.font.color.rgb = RGBColor(*el['color'])
                p.font.bold = el['fontWeight']
                
                if el['textAlign'] == 'center': p.alignment = PP_ALIGN.CENTER
                elif el['textAlign'] == 'right': p.alignment = PP_ALIGN.RIGHT
                else: p.alignment = PP_ALIGN.LEFT
            except Exception:
                pass

    prs.save(output_file)

async def get_html_resolution(browser, html_content: str):
    page = await browser.new_page(viewport={"width": 1280, "height": 720})
    await page.set_content(html_content, wait_until="load")
    await page.evaluate("document.fonts.ready")
    dims = await page.evaluate("""
        () => {
            const preferred = document.querySelector('.slide.active') || document.querySelector('.slide') || document.querySelector('.slide-shell') || document.querySelector('#presentation-container') || document.querySelector('#deck');
            if (preferred) {
                const rect = preferred.getBoundingClientRect();
                if (rect.width > 0 && rect.height > 0) {
                    return { width: Math.ceil(Math.max(rect.width, 1280)), height: Math.ceil(Math.max(rect.height, 720)) };
                }
            }
            const scrollW = Math.max(document.documentElement.scrollWidth, document.body.scrollWidth, document.documentElement.offsetWidth, 1280);
            const scrollH = Math.max(document.documentElement.scrollHeight, document.body.scrollHeight, document.documentElement.offsetHeight, 720);
            return { width: Math.ceil(scrollW), height: Math.ceil(scrollH) };
        }
    """)
    await page.close()
    return dims

async def force_final_state(page, slide_index: int):
    await page.evaluate("""
        (idx) => {
            if (window.goToSlide) { try { window.goToSlide(idx + 1); } catch(e) {} } 
            else if (window.goTo) { try { window.goTo(idx + 1); } catch(e) {} } 
            else {
                const slides = [...document.querySelectorAll('.slide')];
                slides.forEach((s, j) => { if (j === idx) s.classList.add('active'); else s.classList.remove('active'); });
            }
            if (window.runAnims) { try { window.runAnims(idx); } catch(e) {} }
            let slides = [...document.querySelectorAll('.slide')];
            if (slides.length === 0) slides = [...document.querySelectorAll('.slide-shell')];
            if (slides.length > 0) {
                slides.forEach((s, j) => {
                    if (j === idx) {
                        s.classList.add('active');
                        s.style.display = ''; s.style.visibility = 'visible'; s.style.opacity = '1'; s.style.pointerEvents = 'all';
                        s.scrollIntoView({block:'start', inline:'nearest'});
                    } else {
                        s.classList.remove('active');
                        s.style.display = 'none'; s.style.visibility = 'hidden'; s.style.opacity = '0'; s.style.pointerEvents = 'none';
                    }
                });
            }

            const dots = [...document.querySelectorAll('.nav-dots .dot, .dot')];
            dots.forEach((d, j) => {
                if (j === idx) d.classList.add('active');
                else d.classList.remove('active');
            });
            const current = slides.length > 0 ? slides[idx] : document.body;
            if (!current) return;

            const currentCs = getComputedStyle(current);
            if (currentCs.backgroundAttachment === 'fixed') {
                current.style.setProperty('background-attachment', 'scroll', 'important');
            }

            current.querySelectorAll('*').forEach(el => {
                const cs = getComputedStyle(el);
                if (cs.display === 'none') el.style.display = '';
                if (cs.visibility === 'hidden') el.style.visibility = 'visible';
                if (parseFloat(cs.opacity || '1') === 0) el.style.opacity = '1';
                if (cs.backgroundAttachment === 'fixed') {
                    el.style.setProperty('background-attachment', 'scroll', 'important');
                }
                el.style.transition = 'none';
            });
        }
    """, slide_index)

async def render_deck_to_file(html_content: str, output_file: str):
    slide_count = count_slides(html_content)
    slides_for_output = []

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir = Path(temp_dir)
        async with async_playwright() as p:
            browser = await p.chromium.launch()
            dims = await get_html_resolution(browser, html_content)
            context = await browser.new_context(viewport={"width": dims["width"], "height": dims["height"]}, device_scale_factor=2)

            for i in range(slide_count):
                page = await context.new_page()
                await page.set_content(html_content, wait_until="load")
                await page.evaluate("document.fonts.ready") 
                await page.wait_for_timeout(300)

                await page.add_style_tag(content="""
                    * { -webkit-print-color-adjust: exact !important; print-color-adjust: exact !important; background-attachment: scroll !important; }
                    #dots, #ctr, .progress-wrap { display: none !important; }
                    *, *::before, *::after { transition-duration: 0s !important; transition-delay: 0s !important; animation-duration: 0s !important; animation-delay: 0s !important; animation-fill-mode: forwards !important; }
                """)
                
                await force_final_state(page, i)
                await page.wait_for_timeout(2000) 

                await page.evaluate("""
                    () => {
                        if (window.gsap) { try { window.gsap.globalTimeline.progress(1); } catch(e) {} }
                        document.body.querySelectorAll('*').forEach(el => {
                            if (parseFloat(getComputedStyle(el).opacity) === 0) el.style.opacity = '1';
                        });
                    }
                """)

                box = await page.evaluate("""
                    () => {
                        const active = document.querySelector('.slide.active') || document.querySelector('.slide-shell.active') || document.querySelector('.slide') || document.querySelector('.slide-shell') || document.body;
                        const rect = active.getBoundingClientRect();
                        return { x: rect.left, y: rect.top, width: rect.width, height: rect.height };
                    }
                """)
                
                bg_x = max(0, math.floor(box["x"] - 2) if box else 0)
                bg_y = max(0, math.floor(box["y"] - 2) if box else 0)
                bg_w = max(1, math.ceil(box["width"] + 4) if box else dims["width"])
                bg_h = max(1, math.ceil(box["height"] + 4) if box else dims["height"])

                bg_x = min(bg_x, dims["width"] - 1)
                bg_y = min(bg_y, dims["height"] - 1)
                bg_w = min(bg_w, dims["width"] - bg_x)
                bg_h = min(bg_h, dims["height"] - bg_y)
                clip = {"x": bg_x, "y": bg_y, "width": bg_w, "height": bg_h}

                # CRITICAL FIX: Unified Text Engine - Ignores inline elements and extracts full block text
                text_elements = await page.evaluate("""
                    () => {
                        const elements = [];
                        const rgbToHex = (rgba) => {
                            if (!rgba) return null;
                            const match = rgba.match(/^rgba?\\((\\d+),\\s*(\\d+),\\s*(\\d+)/);
                            return match ? [parseInt(match[1]), parseInt(match[2]), parseInt(match[3])] : null;
                        };
                        
                        document.body.querySelectorAll('*').forEach(node => {
                            const style = window.getComputedStyle(node);
                            if (style.display === 'none' || style.visibility === 'hidden' || parseFloat(style.opacity) === 0) return;
                            
                            if (style.display === 'inline') return;

                            const tag = node.tagName.toLowerCase();
                            if (['script', 'style', 'svg', 'i', 'img'].includes(tag)) return;

                            let minX = Infinity, minY = Infinity, maxX = -Infinity, maxY = -Infinity;
                            let hasValidRect = false;

                            const extractTextAndBounds = (element) => {
                                let localText = "";
                                for (let child of element.childNodes) {
                                    if (child.nodeType === 3) { 
                                        const content = child.nodeValue.replace(/\\s+/g, ' '); 
                                        if (content !== ' ' && content !== '') {
                                            localText += content;
                                            const range = document.createRange();
                                            range.selectNode(child);
                                            const rects = range.getClientRects();
                                            for (let r of rects) {
                                                if (r.width > 0 && r.height > 0) {
                                                    minX = Math.min(minX, r.left);
                                                    minY = Math.min(minY, r.top);
                                                    maxX = Math.max(maxX, r.right);
                                                    maxY = Math.max(maxY, r.bottom);
                                                    hasValidRect = true;
                                                }
                                            }
                                        }
                                    } else if (child.nodeType === 1) { 
                                        const childTag = child.tagName.toLowerCase();
                                        if (childTag === 'br') {
                                            localText += '\\n';
                                        } else if (!['script', 'style', 'svg', 'i', 'img'].includes(childTag)) {
                                            const childStyle = window.getComputedStyle(child);
                                            if (childStyle.display === 'inline') {
                                                localText += extractTextAndBounds(child);
                                            }
                                        }
                                    }
                                }
                                return localText;
                            };

                            let textStr = extractTextAndBounds(node).trim();
                            if (!hasValidRect || !textStr) return;

                            const isFlexCenter = style.display === 'flex' && style.alignItems === 'center';

                            elements.push({
                                text: textStr, 
                                x: minX, 
                                y: minY, 
                                w: Math.max(1, maxX - minX + 2), 
                                h: Math.max(1, maxY - minY + 2),
                                fontSize: parseFloat(style.fontSize) || 12, 
                                color: rgbToHex(style.color),
                                fontWeight: style.fontWeight === 'bold' || parseInt(style.fontWeight) >= 600,
                                textAlign: style.textAlign || 'left',
                                vAlign: isFlexCenter ? 'middle' : 'top'
                            });
                        });
                        return elements;
                    }
                """)

                await page.evaluate("""
                    () => {
                        const isCard = (el) => {
                            const tag = el.tagName.toUpperCase();
                            if (['BODY', 'HTML', 'MAIN', 'SECTION', 'HEADER', 'FOOTER', 'SVG', 'SCRIPT', 'STYLE', 'IMG', 'I'].includes(tag)) return false;

                            const excludeClasses = ['slide', 'content-wrap', 'bg-wrap', 'hero-grid', 'hero-split', 'infographic-board', 'outcomes-grid', 'close-layout', 'story-ribbon', 'metric-stack', 'nav', 'top-right'];
                            for (let cls of excludeClasses) { if (el.classList.contains(cls)) return false; }
                            if (el.className && typeof el.className === 'string' && el.className.includes('bg-')) return false;

                            const rect = el.getBoundingClientRect();
                            if (rect.width <= 10 || rect.height <= 10) return false;
                            if (rect.width >= window.innerWidth * 0.85 || rect.height >= window.innerHeight * 0.85) return false;

                            const style = window.getComputedStyle(el);
                            if (style.display === 'none' || style.visibility === 'hidden' || parseFloat(style.opacity) === 0) return false;

                            const explicitClasses = [
                                'ribbon-card', 'metric', 'hero-side', 'issue-card', 'insight-panel',
                                'kpi-box', 'center-hub', 'spoke', 'narrative-panel', 'checkpoint',
                                'outcome-board', 'outcome-row', 'data-source', 'central-engine',
                                'kpi-card', 'phase', 'ey-tag', 'page-chip', 'silo-pillar', 'hub-center', 'counter'
                            ];

                            let isExplicit = false;
                            for (let cls of explicitClasses) { if (el.classList.contains(cls)) isExplicit = true; }
                            if (tag === 'BUTTON') isExplicit = true; 
                            if (isExplicit) return true;

                            const hasBg = (style.backgroundColor !== 'rgba(0, 0, 0, 0)' && style.backgroundColor !== 'transparent') || (style.backgroundImage !== 'none' && style.backgroundImage !== 'initial');
                            const hasBorder = parseFloat(style.borderWidth) > 0 || parseFloat(style.borderTopWidth) > 0;
                            const hasShadow = style.boxShadow !== 'none' && style.boxShadow !== '';
                            return hasBg || hasBorder || hasShadow;
                        };

                        const cards = Array.from(document.body.querySelectorAll('*')).filter(isCard);
                        const topLevelCards = cards.filter(card => {
                            let parent = card.parentElement;
                            while (parent && parent !== document.body) {
                                if (cards.includes(parent)) return false;
                                parent = parent.parentElement;
                            }
                            return true;
                        });

                        topLevelCards.forEach((c, i) => c.setAttribute('data-ppt-card', i));
                    }
                """)

                # FIX: Detect and isolate icon elements within cards
                await page.evaluate("""
                    () => {
                        const isIconElement = (el) => {
                            const tag = el.tagName.toLowerCase();
                            const rect = el.getBoundingClientRect();
                            if (rect.width <= 0 || rect.height <= 0) return false;
                            if (tag === 'i') return true;
                            if (tag === 'svg' && rect.width <= 150 && rect.height <= 150) return true;
                            if (tag === 'img' && rect.width <= 100 && rect.height <= 100) return true;
                            return false;
                        };

                        const isIconContainer = (el) => {
                            if (el.hasAttribute('data-ppt-card')) return false;
                            const rect = el.getBoundingClientRect();
                            if (rect.width < 10 || rect.height < 10 || rect.width > 150 || rect.height > 150) return false;
                            const visibleChildren = Array.from(el.children).filter(c => {
                                const cs = getComputedStyle(c);
                                return cs.display !== 'none' && cs.visibility !== 'hidden' && parseFloat(cs.opacity) > 0;
                            });
                            if (visibleChildren.length === 0) return false;
                            return visibleChildren.every(c => isIconElement(c));
                        };

                        document.querySelectorAll('[data-ppt-card]').forEach(card => {
                            card.querySelectorAll('*').forEach(el => {
                                if (isIconContainer(el)) {
                                    el.setAttribute('data-ppt-icon', '');
                                }
                            });
                            card.querySelectorAll('i, svg').forEach(el => {
                                if (isIconElement(el) && !el.closest('[data-ppt-icon]')) {
                                    el.setAttribute('data-ppt-icon', '');
                                }
                            });
                        });

                        document.querySelectorAll('[data-ppt-icon]').forEach((el, i) => {
                            el.setAttribute('data-ppt-icon', i);
                        });
                    }
                """)

                # CRITICAL FIX: Safe text hiding using TreeWalker to prevent icons from disappearing
                await page.evaluate("""
                    () => {
                        const walker = document.createTreeWalker(document.body, NodeFilter.SHOW_TEXT, null, false);
                        const textNodes = [];
                        
                        while (walker.nextNode()) {
                            const parentTag = walker.currentNode.parentNode.tagName.toLowerCase();
                            if (['script', 'style', 'noscript'].includes(parentTag)) continue;
                            if (walker.currentNode.nodeValue.trim().length > 0) {
                                textNodes.push(walker.currentNode);
                            }
                        }
                        
                        textNodes.forEach(node => {
                            const span = document.createElement('span');
                            span.style.setProperty('color', 'transparent', 'important');
                            span.style.setProperty('-webkit-text-fill-color', 'transparent', 'important');
                            node.parentNode.insertBefore(span, node);
                            span.appendChild(node);
                        });
                    }
                """)

                # Hide icons before capturing card screenshots (preserves layout space)
                await page.evaluate("""
                    () => {
                        document.querySelectorAll('[data-ppt-icon]').forEach(el => {
                            el.style.setProperty('visibility', 'hidden', 'important');
                        });
                    }
                """)

                component_elements = []
                card_count = await page.locator('[data-ppt-card]').count()

                COMPONENT_CLIP_PAD = 8

                for j in range(card_count):
                    loc = page.locator(f'[data-ppt-card="{j}"]')
                    cbox = await loc.bounding_box()
                    
                    if cbox and cbox['width'] > 0 and cbox['height'] > 0:
                        raw_left = cbox['x']
                        raw_top = cbox['y']
                        raw_right = cbox['x'] + cbox['width']
                        raw_bottom = cbox['y'] + cbox['height']

                        left = max(0, raw_left - COMPONENT_CLIP_PAD)
                        top = max(0, raw_top - COMPONENT_CLIP_PAD)
                        right = min(dims['width'], raw_right + COMPONENT_CLIP_PAD)
                        bottom = min(dims['height'], raw_bottom + COMPONENT_CLIP_PAD)

                        cw = right - left
                        ch = bottom - top

                        if cw <= 0 or ch <= 0:
                            continue

                        cx = math.floor(left)
                        cy = math.floor(top)
                        cw = math.ceil(cw)
                        ch = math.ceil(ch)

                        if cx + cw > dims['width']:
                            cw = max(1, dims['width'] - cx)
                        if cy + ch > dims['height']:
                            ch = max(1, dims['height'] - cy)
                        
                        c_clip = {"x": cx, "y": cy, "width": cw, "height": ch}
                        c_img_path = temp_dir / f"comp_{i}_{j}.png"
                        
                        try:
                            await page.screenshot(path=str(c_img_path), clip=c_clip, omit_background=True)
                            component_elements.append({
                                'img_path': c_img_path,
                                'x': cx, 'y': cy, 'w': cw, 'h': ch
                            })
                        except Exception:
                            pass 

                # Restore icons and capture them as separate transparent components
                await page.evaluate("""
                    () => {
                        document.querySelectorAll('[data-ppt-icon]').forEach(el => {
                            el.style.setProperty('visibility', 'visible', 'important');
                        });
                    }
                """)

                ICON_CLIP_PAD = 4
                icon_count = await page.locator('[data-ppt-icon]').count()
                for k in range(icon_count):
                    loc = page.locator(f'[data-ppt-icon="{k}"]')
                    ibox = await loc.bounding_box()

                    if ibox and ibox['width'] > 0 and ibox['height'] > 0:
                        left = max(0, ibox['x'] - ICON_CLIP_PAD)
                        top = max(0, ibox['y'] - ICON_CLIP_PAD)
                        right = min(dims['width'], ibox['x'] + ibox['width'] + ICON_CLIP_PAD)
                        bottom = min(dims['height'], ibox['y'] + ibox['height'] + ICON_CLIP_PAD)

                        iw = right - left
                        ih = bottom - top
                        if iw <= 0 or ih <= 0:
                            continue

                        ix = math.floor(left)
                        iy = math.floor(top)
                        iw = math.ceil(iw)
                        ih = math.ceil(ih)

                        if ix + iw > dims['width']:
                            iw = max(1, dims['width'] - ix)
                        if iy + ih > dims['height']:
                            ih = max(1, dims['height'] - iy)

                        i_clip = {"x": ix, "y": iy, "width": iw, "height": ih}
                        i_img_path = temp_dir / f"icon_{i}_{k}.png"

                        try:
                            await page.screenshot(path=str(i_img_path), clip=i_clip, omit_background=True)
                            component_elements.append({
                                'img_path': i_img_path,
                                'x': ix, 'y': iy, 'w': iw, 'h': ih
                            })
                        except Exception:
                            pass

                await page.evaluate("""
                    () => {
                        document.querySelectorAll('[data-ppt-card]').forEach(el => {
                            el.style.setProperty('visibility', 'hidden', 'important');
                        });
                    }
                """)

                bg_img_path = temp_dir / f"bg_slide_{i+1}.png"
                await page.screenshot(path=str(bg_img_path), clip=clip)

                slides_for_output.append({
                    'bg_img_path': bg_img_path,
                    'width': bg_w, 'height': bg_h,
                    'elements': text_elements, 
                    'components': component_elements,
                    'clip_x': bg_x, 'clip_y': bg_y
                })
                
                await page.close()

            await context.close()
            await browser.close()

        build_layered_pptx(slides_for_output, output_file)

async def generate_pptx(html_content: str, source_filename: str | None = None, original_filename: str | None = None) -> str:
    """Generate PPTX using original uploaded file name when provided."""
    input_name = source_filename or original_filename
    file_name = safe_pptx_filename(input_name)
    output_path = unique_output_path(file_name)
    await render_deck_to_file(html_content, str(output_path))
    return str(output_path)
